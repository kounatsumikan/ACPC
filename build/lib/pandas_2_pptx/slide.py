
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Cm
from pptx.chart.data import CategoryChartData
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

class slide():

    def __init__(self, slide_obj):
        '''
        pptxクラスのslideオブジェクトにデータを追加するクラス

        Args:
            slide (pptx.slide.Slide): pptxクラスのslideオブジェクト。
        '''
        self.__slide = slide_obj
        self.__slide_object_list = []
        self.__df_list = []

    def add_chart(self, df):
        '''
        プライベートメソッド。
        pandasをpptxの棒グラフに変換する関数

        Args:
            df (pandas.DataFrame): 変換するdataframe
        '''
        self.__slide_object_check()
        chart_data = self.__pandas_genelate_linechart(df)
        self.__df_list.append(df)
        self.__slide_object_list.append(["chart", chart_data])

    def add_table(self, df):
        '''
        プライベートメソッド。
        pandasをpptxのテーブルに変換する関数

        Args:
            df (pandas.DataFrame): 変換するdataframe
        '''
        self.__slide_object_check()
        rows = len(df)+1
        cols = len(df.columns)+1
        self.__df_list.append(df)
        self.__slide_object_list.append(["table", rows, cols])

    def __pandas_genelate_table(self, table, df, width):
        '''
        プライベートメソッド。
        pandasをpptxのテーブルに変換する関数

        Args:
            table (pptx.table.Table): tableオブジェクト
            df (pandas.DataFrame): tableオブジェクトに格納するdataframe
        '''
        table.cell(0, 0).text = "index"
        text_frame = table.cell(0, 0).text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        for i, column in enumerate(df.columns):
            run = table.cell(0, i + 1).text_frame.paragraphs[0].add_run()
            font = run.font
            font.name = "メイリオ"
            run.text = column
            run.font.size = width/df.shape[1]

        for i, index in enumerate(df.index):
            table.cell(i + 1, 0).text_frame.paragraphs[0].add_run()
            font = run.font
            font.name = "メイリオ"
            run.text = str(index)
            run.font.size = width/df.shape[1]

        for row in range(df.shape[0]):
            for line in range(df.shape[1]):
                table.cell(row+1, line+1).text_frame.paragraphs[0].add_run()
                font = run.font
                font.name = "メイリオ"
                run.text = str(df.values[row, line])
                run.font.size = width/df.shape[1]

    def __pandas_genelate_linechart(self, df):
        '''
        プライベートメソッド。
        pandasのデータフレームをpptxのchartに変換

        Args:
            df (pandas.DataFrame): 出力したいdataframe

        return:
            chart_data (pptx.chart.data.CategoryChartData): chartデータ
        '''
        chart_data = CategoryChartData()
        chart_data.categories = df.index
        for column in df.columns:
            chart_data.add_series(column, df[column])
        return chart_data

    def __slide_object_check(self):
        '''
        プライベートメソッド。
        グラフやテーブルの数を6つ以上追加しようとするとエラー(それ以上は見栄えが悪くなる)
        '''
        if(len(self.__slide_object_list) == 6):
            raise Exception("objectが多すぎます")

    def __cordinate_shape(self):
        '''
        プライベートメソッド。
        グラフやテーブルの大きさを調整

        return:
            shape_list (list): グラフやテーブルの大きさを格納したリスト
        '''
        object_num = len(self.__slide_object_list)
        split_len = Cm(0.43)
        hight = Cm(14.71)
        width = Cm(24.4)
        left = Inches(0.17)
        top = Inches(1.5)
        shape_list = []

        if(object_num == 1):
            shape_list.append([left + split_len * 2,
                               top + split_len*2,
                               Cm(23),
                               Cm(13)])

        if(object_num <= 2):
            width = int((width - split_len) / 2)
            hight = int((hight - split_len*2) / 1)
            for i in range(object_num):
                shape_list.append([left + (width + split_len) * i, top, width, hight])

        elif(object_num <= 4):
            width = int((width - split_len) / 2)
            hight = int((hight - split_len*2) / 2)
            for i in range(2):
                for j in range(2):
                    shape_list.append([left + (width + split_len) * j, top + (hight + split_len) * i, width, hight])

        elif(object_num <= 6):
            width = int((width - split_len * 2) / 3)
            hight = int((hight - split_len*2) / 2)
            for i in range(2):
                for j in range(3):
                    shape_list.append([left + (width + split_len) * j, top + (hight + split_len) * i, width, hight])
        return shape_list

    def generate(self):
        '''
        プライベートメソッド。
        追加したグラフやテーブルを__cordinate_shape()の情報をもとにスライドに書き込む
        '''
        shape_list = self.__cordinate_shape()
        for n, i in enumerate(self.__slide_object_list):
            if(i[0] == "table"):
                table = self.__slide.shapes.add_table(i[1],  # row
                                                      i[2],  # col
                                                      shape_list[n][0],  # left
                                                      shape_list[n][1],  # top
                                                      shape_list[n][2],  # width
                                                      shape_list[n][3]   # hight
                                                      ).table
                self.__pandas_genelate_table(table, self.__df_list[n], shape_list[n][2])

            elif(i[0] == "chart"):
                chart = self.__slide.shapes.add_chart(XL_CHART_TYPE.LINE,
                                                      shape_list[n][0],  # left
                                                      shape_list[n][1],  # top
                                                      shape_list[n][2],  # width
                                                      shape_list[n][3],  # hight
                                                      i[1]
                                                      ).chart
                chart.has_legend = True
                chart.legend.include_in_layout = False
