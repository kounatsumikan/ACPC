
from pptx import Presentation
import sys
import pathlib
# base.pyのあるディレクトリの絶対パスを取得
current_dir = pathlib.Path(__file__).resolve().parent
# モジュールのあるパスを追加
sys.path.append(str(current_dir) + '/../')
from .slide import slide


class pandas_2_pptx():

    def __init__(self, fname=None):
        '''
        pandasのdataframeをpptxのスライドとして出力するクラス

        Args:
            fname (string): pptxファイルのパス。指定した場合、
            読み込んだpptxにデータが追記される。指定しない場合新しいファイルに出力される。
        '''
        if(fname is None):
            self.prs = Presentation()
        else:
            self.prs = Presentation(fname)
            print(f"{fname}を読み込みました。\n上書きに気をつけてください")

    def add_title_slide(self, title_txt):
        '''
        タイトルスライドを追加する

        Args:
            title_txt (string): スライドのタイトル名

        return:
            slide (slide): slideオブジェクト
        '''
        slide_obj = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        title_slide = slide_obj.shapes.title
        title_slide.text = title_txt
        slide_obj = slide(slide_obj)
        return slide_obj

    def add_slide(self, title_txt):
        '''
        スライドを追加する

        Args:
            title_txt (string): スライドのタイトル名

        return:
            slide (slide): slideオブジェクト
        '''
        slide_obj = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        title = slide_obj.shapes.title
        title.text = title_txt
        slide_obj = slide(slide_obj)
        return slide_obj

    def save(self, fname):
        '''
        スライドを追加したpptxデータを出力する

        Args:
            fname (string): 出力するpptxファイルのパス。
        '''
        self.prs.save(fname)
